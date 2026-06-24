"""Convert RLimage HTML presentation to PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG_DEEP = RGBColor(0x0A, 0x0F, 0x1A)
BG_NAVY = RGBColor(0x11, 0x1B, 0x2E)
BG_SLATE = RGBColor(0x1A, 0x2A, 0x42)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
GOLD_LIGHT = RGBColor(0xE8, 0xC8, 0x7A)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
TEXT_MAIN = RGBColor(0xF0, 0xEC, 0xE4)
TEXT_SUB = RGBColor(0x9F, 0xB3, 0xC8)
TEXT_WEAK = RGBColor(0x5A, 0x70, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0x25, 0x2E, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # blank layout

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_MAIN, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='SimSun', italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.font.name = font_name
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return tb

def add_multiline_textbox(slide, left, top, width, height, lines, line_height=1.3):
    """lines is list of (text, font_size, color, bold, alignment, font_name) tuples"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text, font_size, color, bold, alignment, font_name = line_data
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.font.name = font_name
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return tb

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if border_color:
        rect.line.color.rgb = border_color
        rect.line.width = Pt(0.5)
    else:
        rect.line.fill.background()
    return rect

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if border_color:
        rect.line.color.rgb = border_color
        rect.line.width = Pt(0.5)
    else:
        rect.line.fill.background()
    return rect

# ============================================================
# SLIDE 1: COVER
# ============================================================
s1 = prs.slides.add_slide(BLANK)
set_bg(s1, BG_DEEP)

# Decor ring (approximated as two unfilled circles)
c1 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(0.5), Inches(1.5), Inches(1.5))
c1.fill.background()
c1.line.color.rgb = GOLD
c1.line.width = Pt(1.5)
c1.line.dash_style = 1  # solid

c2 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(0.7), Inches(1.1), Inches(1.1))
c2.fill.background()
c2.line.color.rgb = GOLD
c2.line.width = Pt(0.75)

# Tag
add_textbox(s1, 0, 1.3, 13.333, 0.5, "RESEARCH REPORT", font_size=14, color=GOLD,
            alignment=PP_ALIGN.CENTER, font_name='Georgia', bold=True)

# Title
add_textbox(s1, 1, 2.0, 11.333, 1.8, "面向视觉感知模型的\n可验证奖励后训练", font_size=42, color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')

# Divider
add_rect(s1, 5.5, 3.9, 2.3, 0.04, GOLD)

# Subtitle
add_textbox(s1, 0, 4.2, 13.333, 0.6, "RLimage: Verifiable Reward Vision Post-Training",
            font_size=18, color=TEXT_SUB, alignment=PP_ALIGN.CENTER)

# Info row
for i, (label, value) in enumerate([("项目", "RLimage"), ("日期", "2026-06-04"), ("版本", "v1.0")]):
    x = 3.5 + i * 2.5
    add_textbox(s1, x, 5.6, 2.0, 0.3, label, font_size=12, color=TEXT_WEAK, alignment=PP_ALIGN.CENTER)
    add_textbox(s1, x, 5.9, 2.0, 0.4, value, font_size=16, color=TEXT_SUB, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: BACKGROUND + TECH ROUTE (Split)
# ============================================================
s2 = prs.slides.add_slide(BLANK)
set_bg(s2, BG_NAVY)

# Left panel
add_textbox(s2, 1, 0.3, 3, 1.0, "01", font_size=64, color=GOLD, bold=True, font_name='Georgia')
add_textbox(s2, 1, 1.2, 5, 0.8, "研究背景", font_size=28, color=WHITE, bold=True, font_name='SimHei')

# Card 1 - gold left border
add_rounded_rect(s2, 1, 2.2, 5.2, 2.4, BG_SLATE, CARD_BORDER)
add_rect(s2, 1, 2.25, 0.06, 2.3, TEAL)  # left accent
add_textbox(s2, 1.4, 2.3, 4.6, 0.5, "核心问题", font_size=16, color=TEAL, bold=True, font_name='SimHei')
add_textbox(s2, 1.4, 2.8, 4.6, 1.6,
    "将文本模型领域的 RLVR（可验证奖励强化学习）迁移到视觉目标检测，使 detector 通过频域证据学会区分可信候选框",
    font_size=15, color=TEXT_SUB)

# Card 2 - gold left border
add_rounded_rect(s2, 1, 4.8, 5.2, 1.8, BG_SLATE, CARD_BORDER)
add_rect(s2, 1, 4.85, 0.06, 1.7, TEAL)
add_textbox(s2, 1.4, 4.9, 4.6, 0.5, "关键挑战", font_size=16, color=TEAL, bold=True, font_name='SimHei')
add_textbox(s2, 1.4, 5.4, 4.6, 1.0,
    "输出是候选框集合（类别+坐标+NMS），频域 verifier 需绑定到具体 ROI action，而非图像级 loss 缩放",
    font_size=15, color=TEXT_SUB)

# Right panel
add_textbox(s2, 7, 0.3, 3, 1.0, "02", font_size=64, color=GOLD, bold=True, font_name='Georgia')
add_textbox(s2, 7, 1.2, 5, 0.8, "技术路线", font_size=28, color=WHITE, bold=True, font_name='SimHei')

add_rounded_rect(s2, 7, 2.2, 5.2, 2.4, BG_SLATE, CARD_BORDER)
add_rect(s2, 7, 2.25, 0.06, 2.3, GOLD)
add_textbox(s2, 7.4, 2.3, 4.6, 0.5, "五阶段递进", font_size=16, color=GOLD, bold=True, font_name='SimHei')
add_textbox(s2, 7.4, 2.8, 4.6, 1.6,
    "MVP 频域诊断 → Spectral Quality Head → RLVR 后训练 → MPLSeg AFM 网络内 FFT → 语义分割迁移",
    font_size=15, color=TEXT_SUB)

add_rounded_rect(s2, 7, 4.8, 5.2, 1.8, BG_SLATE, CARD_BORDER)
add_rect(s2, 7, 4.85, 0.06, 1.7, TEAL)
add_textbox(s2, 7.4, 4.9, 4.6, 0.5, "基础设施", font_size=16, color=TEAL, bold=True, font_name='SimHei')
add_textbox(s2, 7.4, 5.4, 4.6, 1.0,
    "Penn-Fudan / VOC 数据集 · Faster R-CNN · spectral verifier · signed advantage · KL anchor · shuffled control",
    font_size=14, color=TEXT_SUB)


# ============================================================
# SLIDE 3: CORE WORK (3 Cards)
# ============================================================
s3 = prs.slides.add_slide(BLANK)
set_bg(s3, BG_DEEP)

add_textbox(s3, 0, 0.6, 13.333, 0.8, "核心工作", font_size=32, color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')
add_textbox(s3, 0, 1.3, 13.333, 0.4, "CORE CONTRIBUTIONS", font_size=13, color=TEXT_WEAK,
            alignment=PP_ALIGN.CENTER, font_name='Georgia')

cards = [
    {
        'x': 0.8, 'top_border': GOLD, 'icon_char': '🔬', 'icon_color': GOLD,
        'title': 'RLVR 稳定 Shell',
        'desc': '实现 rollout → verifier → signed advantage → KL anchor 完整 pipeline。KL=10 + policy=0.0003 使 AP50 与 baseline 差距 < 0.01，首次在检测任务中稳定运行 RLVR 后训练。'
    },
    {
        'x': 4.6, 'top_border': TEAL, 'icon_char': '📊', 'icon_color': TEAL,
        'title': '频域因果检验',
        'desc': '设计 real vs shuffled spectral + IoU-only + structure 等多组对照。R_amp TP/FP AUC 达 0.93，但 post-training 中 real 与 shuffled 无显著差异，手工频域 verifier 因果信号不足。'
    },
    {
        'x': 8.4, 'top_border': GOLD, 'icon_char': '⚡', 'icon_color': GOLD,
        'title': 'MPLSeg AFM\n梯度修复',
        'desc': '修复旧 AFM 的梯度死锁：mag_scale 和 phase_scale 始终为 0。MPLSeg 式硬激活门控使 residual_scale=0.972，FFT 路径被真正激活，8 项单元测试全部通过。'
    }
]

for card in cards:
    x = card['x']
    # Card background
    add_rounded_rect(s3, x, 2.1, 3.6, 4.5, BG_SLATE, CARD_BORDER)
    # Top border
    add_rect(s3, x + 0.2, 2.1, 3.2, 0.04, card['top_border'])
    # Icon circle
    circle = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.4), Inches(2.5), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    if card['icon_color'] == GOLD:
        circle.fill.fore_color.rgb = RGBColor(0x3D, 0x32, 0x18)
    else:
        circle.fill.fore_color.rgb = RGBColor(0x16, 0x3E, 0x38)
    circle.line.fill.background()
    # Card title
    title_lines = card['title'].split('\n')
    add_textbox(s3, x + 0.4, 3.4, 2.8, 0.7, title_lines[0], font_size=18, color=WHITE,
                bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')
    if len(title_lines) > 1:
        add_textbox(s3, x + 0.4, 3.8, 2.8, 0.5, title_lines[1], font_size=18, color=WHITE,
                    bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')
        desc_top = 4.2
    else:
        desc_top = 3.9
    add_textbox(s3, x + 0.3, desc_top, 3.0, 2.2, card['desc'], font_size=14, color=TEXT_SUB,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: RESULTS (Data Viz)
# ============================================================
s4 = prs.slides.add_slide(BLANK)
set_bg(s4, BG_NAVY)

# Left panel
add_textbox(s4, 1, 0.6, 6, 0.8, "关键实验结果", font_size=30, color=WHITE, bold=True, font_name='SimHei')
add_textbox(s4, 1, 1.3, 6, 0.8, "Penn-Fudan + VOC 数据集，14 轮实验迭代，从分类假设到分割迁移的完整证据链",
            font_size=15, color=TEXT_SUB)

metrics = [
    ("0.885", "RLVR Stable Shell AP50"),
    ("0.956", "QH ROI-only TP/FP AUC"),
    ("0.972", "MPLSeg AFM residual_scale"),
    ("14", "实验轮次完整记录"),
]
for i, (num, label) in enumerate(metrics):
    y = 2.5 + i * 1.1
    add_textbox(s4, 1, y, 1.8, 0.8, num, font_size=36, color=GOLD, bold=True, font_name='Georgia')
    add_textbox(s4, 2.8, y + 0.15, 4.5, 0.5, label, font_size=14, color=TEXT_WEAK)

# Right panel: bar chart
bar_data = [
    ("RLVR\nAP50", 0.885, GOLD, 230),
    ("MPLSeg\nAP50", 0.868, GOLD, 225),
    ("VOC\nAP50", 0.774, TEAL, 200),
    ("AP75", 0.653, GOLD, 170),
    ("ECE↓", 0.067, TEAL, 100),
]

for i, (label, val, color, h_px) in enumerate(bar_data):
    bx = 8.0 + i * 0.98
    bar_top = 6.7 - h_px * 0.011  # scale to inches
    bar_h = h_px * 0.011

    # Bar value label
    add_textbox(s4, bx - 0.1, bar_top - 0.4, 0.9, 0.35, str(val), font_size=14, color=TEXT_SUB,
                alignment=PP_ALIGN.CENTER, bold=True)

    # Bar
    bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(bar_top), Inches(0.6), Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Bar label
    add_textbox(s4, bx - 0.15, 6.85, 0.9, 0.6, label, font_size=14, color=TEXT_WEAK,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: CONCLUSION
# ============================================================
s5 = prs.slides.add_slide(BLANK)
set_bg(s5, BG_DEEP)

# Glow (approximated)
glow = s5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1.5), Inches(5.3), Inches(5.3))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x2E)
glow.line.fill.background()

add_textbox(s5, 0, 1.8, 13.333, 0.5, "CONCLUSION", font_size=13, color=GOLD,
            alignment=PP_ALIGN.CENTER, font_name='Georgia', bold=True)
add_textbox(s5, 0, 2.3, 13.333, 1.0, "核心结论", font_size=52, color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')
add_textbox(s5, 0, 3.3, 13.333, 0.6,
            "检测 RLVR shell 已稳定 · 手工频域 verifier 因果信号不足 · 下一站：语义分割",
            font_size=18, color=TEXT_SUB, alignment=PP_ALIGN.CENTER, italic=True)

# Divider
add_rect(s5, 5.8, 4.1, 1.6, 0.03, GOLD)

# Info text
add_textbox(s5, 1.5, 4.5, 10.333, 2.5,
    "视觉检测 RLVR 后训练框架可以稳定实现\n手工 ROI 傅里叶 verifier 在 Penn-Fudan/VOC 上没有因果证据\nMPLSeg-style AFM 修复了梯度死锁，但 box-level frequency reward 仍未成立\n\n→ 迁移至语义分割（dense mask + 幅度/相位解耦原生契合）",
    font_size=15, color=TEXT_SUB, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = "/mnt/c/Users/青云志/Desktop/RLimage_Research_Report.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
