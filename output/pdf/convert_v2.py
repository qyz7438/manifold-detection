"""Convert RLimage v2 presentation (14 slides) to PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG_DEEP = RGBColor(0x0A, 0x0F, 0x1A)
BG_NAVY = RGBColor(0x11, 0x1B, 0x2E)
BG_SLATE = RGBColor(0x1A, 0x2A, 0x42)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
RED = RGBColor(0xE0, 0x55, 0x6A)
GREEN = RGBColor(0x4E, 0xCB, 0x71)
TEXT_MAIN = RGBColor(0xF0, 0xEC, 0xE4)
TEXT_SUB = RGBColor(0x9F, 0xB3, 0xC8)
TEXT_WEAK = RGBColor(0x5A, 0x70, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0x1A, 0x2A, 0x42)
CARD_BORDER = RGBColor(0x25, 0x2E, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, font_size=18, color=TEXT_MAIN,
       bold=False, alignment=PP_ALIGN.LEFT, font_name='SimSun', italic=False):
    """add textbox helper"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.font.name = font_name
    return box

def rect(slide, left, top, width, height, fill_color, border_color=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    r.fill.solid()
    r.fill.fore_color.rgb = fill_color
    if border_color:
        r.line.color.rgb = border_color
        r.line.width = Pt(0.5)
    else:
        r.line.fill.background()
    return r

def plan_slide(slide, badge_text, badge_bg, badge_color, title, desc,
               metrics, verdict_text, verdict_type):
    """Standard plan slide layout.
    metrics: list of (num, num_color, label, name)
    verdict_type: 'pos'/'neg'/'mix'/'info'
    """
    set_bg(slide, BG_NAVY)

    # Left column
    # Badge
    badge = rect(slide, 0.8, 0.6, 2.2, 0.42, badge_bg, rounded=True)
    tb(slide, 0.9, 0.62, 2.0, 0.38, badge_text, font_size=14, color=badge_color,
       bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    title_lines = title.split('\n')
    y_title = 1.2
    for line in title_lines:
        tb(slide, 0.8, y_title, 5, 0.6, line, font_size=26, color=WHITE, bold=True, font_name='SimHei')
        y_title += 0.6

    # Description
    tb(slide, 0.8, y_title + 0.2, 4.5, 1.5, desc, font_size=14, color=TEXT_SUB)

    # Right column: metric cards
    y_m = 0.75
    for num, num_color, label, name in metrics:
        card = rect(slide, 5.8, y_m, 6.8, 0.95, CARD_BG, CARD_BORDER, rounded=True)
        tb(slide, 6.0, y_m + 0.1, 1.3, 0.7, num, font_size=26, color=num_color,
           bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
        tb(slide, 7.4, y_m + 0.05, 5.0, 0.35, label, font_size=14, color=TEXT_WEAK)
        tb(slide, 7.4, y_m + 0.42, 5.0, 0.4, name, font_size=14, color=TEXT_SUB)
        y_m += 1.1

    # Verdict
    verd_colors = {
        'pos': (GREEN, RGBColor(0x0D, 0x2D, 0x1A)),
        'neg': (RED, RGBColor(0x2D, 0x0F, 0x14)),
        'mix': (GOLD, RGBColor(0x25, 0x20, 0x12)),
        'info': (TEAL, RGBColor(0x0E, 0x2A, 0x26)),
    }
    vc, vb = verd_colors.get(verdict_type, (TEXT_SUB, CARD_BG))
    vbox = rect(slide, 5.8, y_m + 0.1, 6.8, 0.55, vb, rounded=True)
    rect(slide, 5.8, y_m + 0.1, 0.05, 0.55, vc)
    tb(slide, 6.1, y_m + 0.18, 6.3, 0.5, verdict_text, font_size=14, color=vc)


# ================================================================
# SLIDE 1: COVER
# ================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DEEP)
tb(s, 0, 1.3, 13.333, 0.5, "RLIMAGE · 完整实验迭代", font_size=14, color=GOLD,
   bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
tb(s, 1, 2.0, 11.333, 1.8, "面向视觉感知模型的\n可验证奖励后训练", font_size=40, color=WHITE,
   bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')
rect(s, 5.5, 3.9, 2.3, 0.04, GOLD)
tb(s, 0, 4.2, 13.333, 0.6, "Verifiable Reward Vision Post-Training — 14 轮实验完整记录",
   font_size=17, color=TEXT_SUB, alignment=PP_ALIGN.CENTER)
for i, (label, value) in enumerate([("项目", "RLimage"), ("日期", "2026-06-04"), ("版本", "v1.0")]):
    x = 3.5 + i * 2.5
    tb(s, x, 5.6, 2.0, 0.3, label, font_size=14, color=TEXT_WEAK, alignment=PP_ALIGN.CENTER)
    tb(s, x, 5.9, 2.0, 0.4, value, font_size=16, color=TEXT_SUB, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2: MVP
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 1", RGBColor(0x0E,0x2A,0x26), TEAL,
           "MVP\nR_amp 初步证据",
           "Penn-Fudan person detection baseline + ROI 幅度谱相似度 R_amp 诊断 TP vs FP。image-level reward-weighted fine-tuning 测试。",
           [("0.930", TEAL, "R_amp TP/FP AUC", "频域证据存在统计判别信号"),
            ("0.863", RED, "Baseline → Posttrain AP50", "Reward posttrain 崩至 0.435"),
            ("0.462", RED, "Posttrain Recall", "image-level reward 导致极度保守")],
           "⚠ R_amp 排序 AUC 高，但直接乘 detector loss 导致模型崩盘。诊断信号 ≠ 训练信号。",
           'neg')

# ================================================================
# SLIDE 3: Quality Head
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 2", RGBColor(0x25,0x20,0x12), GOLD,
           "Spectral\nQuality Head",
           "冻结 detector，训练质量头预测 ROI 质量 q = class_correct × IoU × R_amp。候选框重排序替代 loss 乘法。",
           [("0.956", TEAL, "QH ROI-only TP/FP AUC", "远超 MVP 的手工 R_amp"),
            ("0.874", TEAL, "Rerank α=0.9 AP50", "ECE 0.166, high-conf FP 减半"),
            ("0.947", RED, "ROI+Amp+Struct AUC", "频域特征未超过 ROI-only")],
           "✓ 校准/重排序有效 · ✗ 频域无增量 · ✗ 仍是 reranking 非 RLVR",
           'mix')

# ================================================================
# SLIDE 4: NNI
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 3", RGBColor(0x25,0x20,0x12), GOLD,
           "NNI\nQuality Matrix",
           "108 trials 全面消融: detector epochs (1/3/5) × QH type × QH epochs (8/20) × α (0.7-0.95)",
           [("108", TEAL, "总 Trial 数（全部成功）", "多变量完整矩阵"),
            ("0.860", GOLD, "最佳 AP50", "Det-1ep QH-ROI α=0.70, ECE 0.066"),
            ("0.876", GOLD, "Precision@R=0.85", "α=0.7 激进校准 vs α=0.9 保守")],
           "关键结论: ① α是真 trade-off ② 频域从未超过 ROI-only ③ QH 最大收益是校准非 AP",
           'mix')

# ================================================================
# SLIDE 5: RLVR R1→2.3
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 4-7", RGBColor(0x2D,0x0F,0x14), RED,
           "RLVR 后训练\n从崩塌到稳定",
           "R1: NNI search → R2.1: 工程修复 → R2.2: signed advantage → R2.3: stable shell 成立",
           [("0.668", RED, "Round 1 NNI 最佳 AP50", "ECE 降至 0.087 但 AP 代价大"),
            ("0.623", RED, "R2.1 修复后 AP50", "低置信预测爆炸, precision 0.21"),
            ("0.873", GREEN, "R2.3 Stable Shell AP50", "KL=10, policy=0.0003, frozen baseline")],
           "✓ RLVR shell 首次在检测任务中稳定: signed IoU + KL anchor + frozen baseline, AP 差距 <0.01",
           'pos')

# ================================================================
# SLIDE 6: Round 2.5
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 8", RGBColor(0x2D,0x0F,0x14), RED,
           "Round 2.5\n频谱因果检验",
           "补齐 real vs shuffled spectral + IoU-only + structure 等对照。验证手工频域 verifier 因果性。",
           [("0.885", GOLD, "signed_amp AP50", "8组全部 0.884-0.886"),
            ("0.886", RED, "shuffled control AP50", "与 real 差距 < 0.001"),
            ("0.664", GOLD, "amp AP75 vs null 0.644", "轻微提升但 structure 也类似")],
           "⚠ 核心负结论: 手工 box-level 频谱 verifier 在 Penn-Fudan 上没有因果信号",
           'neg')

# ================================================================
# SLIDE 7: MPLSeg AFM R2.6-2.8
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 9-11", RGBColor(0x0E,0x2A,0x26), TEAL,
           "MPLSeg AFM\n网络内 FFT",
           "将幅度/相位解耦操作移入网络内部: MicroAFM 模块。R2.6 普通版 → R2.7 identity 残差 → R2.8 诊断。",
           [("0.853", RED, "R2.6 AFM Full AP50", "AP75 0.409, precision 0.439"),
            ("0.876", GOLD, "R2.7 Identity AFM AP50", "保住 AP50 但 AP75 仍降"),
            ("0.738", RED, "R2.8 AFM+Box AP75", "mag_scale=0, phase_scale=0: gate 未激活!")],
           "⚠ FFT gate 从未激活，AP75 提升来自 box_head 适配而非频域门控",
           'neg')

# ================================================================
# SLIDE 8: R2.9-2.10
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 12", RGBColor(0x25,0x20,0x12), GOLD,
           "R2.9-2.10\n后训练补洞",
           "修复 checkpoint loading 与 edge-mix 评估。保留必要后训练维度，不再重复频谱无效实验。",
           [("0.815", GOLD, "B2 Posttrain AP50", "AP75 0.525, ECE 0.070"),
            ("0.735", TEAL, "Posttrain Precision", "high-conf FP 5→3, precision 提升"),
            ("0.871", GOLD, "RPN Mixed AP50", "ECE 0.045, recall 0.890")],
           "✓ 后训练稳定运行，precision/ECE 改善 · ✗ AP75 可能降 · ✗ 非频域因果",
           'mix')

# ================================================================
# SLIDE 9: Plan 2.11 VOC
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 13", RGBColor(0x2D,0x0F,0x14), RED,
           "Plan 2.11\nVOC Signal Gate",
           "'是不是任务太简单?' — 换 VOC 3 类检测子集验证手工频谱 verifier 是否超过 shuffled control。",
           [("0.772", GOLD, "V4 Spatial+Spectral AP50", "prec 0.305, recall 0.855"),
            ("0.774", RED, "V5 Shuffled Spectral AP50", "shuffled 略优 real!"),
            ("0.773", GOLD, "V2 Detection-only AP50", "无频谱分支的基线")],
           "⚠ V4 未超 V3/V5: 更难检测子集也没有救回手工频谱 verifier。不再继续 VOC/COCO 检测后训练。",
           'neg')

# ================================================================
# SLIDE 10: Plan 2.12
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "ROUND 14", RGBColor(0x0D,0x2D,0x1A), GREEN,
           "Plan 2.12\nMPLSeg AFM\n梯度修复",
           "针对 R2.8 暴露的架构根因: mag_scale/phase_scale 严格零梯度。用硬激活门控替代可学习 scale。8 项单元测试。",
           [("0.868", GREEN, "MPLSeg AFM AP50", "vs identity AFM 0.863"),
            ("0.653", GREEN, "MPLSeg AFM AP75", "vs identity 0.438 (Δ+0.215!)"),
            ("0.972", TEAL, "residual_scale", "FFT 路径被真正激活!")],
           "✓ FFT 路径恢复梯度 · ✓ residual_scale 非零 · ✓ 单元测试 8/8 · → 但≠ box-level reward 有效",
           'pos')

# ================================================================
# SLIDE 11: Plan 3.x
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "讨论", RGBColor(0x25,0x20,0x12), GOLD,
           "Plan 3.x\n大规模搜索为何不做",
           "Phase A 68 + Phase B 54 + Phase C 24 = 146 trials。原计划判断 real amp 是否优于 shuffled amp。",
           [("146", GOLD, "原计划总 Trial 数", "三阶段完整超参搜索"),
            ("<0.001", RED, "R2.5 8组 AP50 差距", "继续搜参边际收益极低"),
            ("3", GOLD, "替代路线", "换更大数据集 / 学习式 verifier / 写负结果论文")],
           "决策: 不盲目大搜。先写报告识别缺口，再决定 VOC/COCO 或 segmentation。",
           'mix')

# ================================================================
# SLIDE 12: Plan 4.x
# ================================================================
s = prs.slides.add_slide(BLANK)
plan_slide(s, "NEXT", RGBColor(0x0E,0x2A,0x26), TEAL,
           "Plan 4.x\n迁移到语义分割",
           "新包 spectral_segmentation_posttrain。Dense mask 为幅度/相位 verifier 提供像素级可验证目标。7 组对照实验。",
           [("7", TEAL, "实验组数 S1-S7", "baseline→spatial→amp→structure→full→shuffled"),
            ("Dice/IoU/BF1", GOLD, "Dense Verifier", "Mask-level 可验证奖励"),
            ("S6 > S3", GOLD, "Promotion Rule", "S6 超 S3 且 S7 ≠ S6 才声称因果")],
           "→ 下一步主线: Penn-Fudan binary mask RLVR + amplitude/phase/spatial verifier",
           'info')

# ================================================================
# SLIDE 13: Summary Table
# ================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DEEP)
tb(s, 0, 0.4, 13.333, 0.8, "14 轮迭代全景", font_size=30, color=WHITE,
   bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')

data = [
    ["1", "MVP", "R_amp reward posttrain", "0.863→0.435", "—", "直接乘 loss 崩 recall"],
    ["2", "QH", "Spectral Quality Head", "0.874", "—", "AUC 0.956, 校准有效"],
    ["3", "NNI", "108 trials 消融", "0.860", "—", "ROI-only 最优, 频域无增量"],
    ["4-7", "RLVR R1-2.3", "signed advantage + KL", "0.873", "0.658", "Stable shell 成立"],
    ["8", "R2.5 Causality", "real vs shuffled", "0.885", "0.664", "real≈shuffled, 因果不成立"],
    ["9-11", "AFM R2.6-8", "In-network FFT gate", "0.865", "0.738", "Gate 未激活"],
    ["12", "R2.10", "Checkpoint修复", "0.815", "0.525", "后训练稳定但非频域因果"],
    ["13", "P2.11 VOC", "3类检测 signal gate", "0.772", "0.375", "shuffled > real"],
    ["14", "P2.12 AFM Fix", "MPLSeg硬激活门控", "0.868", "0.653", "FFT路径恢复, resid=0.97"],
]

n_cols = 6
col_widths = [0.6, 1.7, 2.2, 1.6, 1.0, 2.1]
col_x = [0.6]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# Header
for ci, (header, x, w) in enumerate(zip(["#","阶段","核心方法","AP50","AP75","关键发现"], col_x, col_widths)):
    tb(s, x, 1.3, w, 0.4, header, font_size=14, color=GOLD, bold=True, font_name='SimHei')
rect(s, 0.5, 1.75, 10.6, 0.02, RGBColor(0x30,0x28,0x18))

# Data rows
y = 1.9
for row in data:
    bg = RGBColor(0x12, 0x18, 0x26) if data.index(row) % 2 == 0 else BG_DEEP
    rect(s, 0.5, y - 0.05, 10.7, 0.45, bg)
    for ci, (val, x, w) in enumerate(zip(row, col_x, col_widths)):
        tb(s, x, y, w, 0.35, val, font_size=14, color=TEXT_SUB)
    y += 0.45

# Bottom row
rect(s, 0.5, y + 0.05, 10.7, 0.02, TEAL)
tb(s, 0.6, y + 0.2, 10.5, 0.4, "→ 下一站: 语义分割 Plan 4.x — Dense mask RLVR",
   font_size=14, color=TEAL, bold=True, font_name='SimHei')

# ================================================================
# SLIDE 14: Conclusion
# ================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_DEEP)

# Glow
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1.5), Inches(5.3), Inches(5.3))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x2E)
glow.line.fill.background()

tb(s, 0, 1.8, 13.333, 0.5, "CONCLUSION", font_size=14, color=GOLD,
   bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
tb(s, 0, 2.3, 13.333, 1.0, "核心结论", font_size=48, color=WHITE,
   bold=True, alignment=PP_ALIGN.CENTER, font_name='SimHei')
tb(s, 0, 3.2, 13.333, 0.6,
   "检测 RLVR shell 已稳定 · 手工频域 verifier 因果不足 · 下一站: 语义分割",
   font_size=17, color=TEXT_SUB, alignment=PP_ALIGN.CENTER, italic=True)
rect(s, 5.8, 4.0, 1.6, 0.03, GOLD)
tb(s, 1.5, 4.3, 10.333, 2.5,
   "① RLVR 后训练 shell 在目标检测上首次稳定实现（14轮迭代）\n"
   "② 手工 ROI 傅里叶 verifier: 统计信号有（AUC 0.93-0.96），因果信号无（real≈shuffled）\n"
   "③ MPLSeg-style AFM 修复了梯度死锁，但≠ box-level frequency reward 有效\n"
   "④ 语义分割为幅度/相位解耦提供原生契合的 dense mask 空间",
   font_size=15, color=TEXT_SUB, alignment=PP_ALIGN.CENTER)

# ================================================================
# SAVE
# ================================================================
output_path = "/mnt/c/Users/青云志/Desktop/RLimage_Research_Report_v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
